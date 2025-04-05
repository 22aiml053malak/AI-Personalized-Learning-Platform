import os
import json
import requests
import logging
from functools import wraps
from flask import Flask, render_template, redirect, url_for, flash, request, jsonify
from flask_login import login_user, current_user, logout_user, login_required, UserMixin
from config import Config
from ext import db, bcrypt, login_manager
from quiz_generator import generate_quiz  # ✅ Import Quiz Generator
import ollama  # ✅ Import Ollama for AI processing
from flask_migrate import Migrate
import random
# ✅ Flask App Initialization
from flask import Flask, render_template, redirect, url_for, flash
from flask_login import login_user, logout_user, login_required, current_user
from flask_migrate import Migrate
from functools import wraps

from config import Config
from ext import db, bcrypt, login_manager  # ✅ Importing extensions from ext.py
from models import User, Course, Enrollment  # ✅ Importing models from model.py

# ✅ Initialize Flask App
app = Flask(__name__)


app.config.from_object(Config)

# ✅ Initialize Flask Extensions
db.init_app(app)
bcrypt.init_app(app)
login_manager.init_app(app)
migrate = Migrate(app, db)

# ✅ Fix SQLAlchemy 2.0 Warning in User Loader
@login_manager.user_loader
def load_user(user_id):
    return db.session.get(User, int(user_id))  # ✅ Fixed for SQLAlchemy 2.0

# ✅ Create database tables at startup
with app.app_context():
    db.create_all()

# ✅ Inject user into all templates
@app.context_processor
def inject_user():
    return dict(user=current_user)

# ✅ Admin Role Decorator
def admin_required(f):
    @wraps(f)
    def decorated_function(*args, **kwargs):
        if not current_user.is_authenticated or not current_user.is_admin:
            flash("You are not authorized to access this page", "danger")
            return redirect(url_for('home'))
        return f(*args, **kwargs)
    return decorated_function

# ✅ Home Page
@app.route("/")
def home():
    courses = Course.query.all()
    return render_template('home.html', courses=courses)

# ✅ User Registration Route
@app.route("/register", methods=['GET', 'POST'])
def register():
    if current_user.is_authenticated:
        return redirect(url_for('home'))
    from forms import RegistrationForm  # ✅ Import inside function to avoid circular import
    form = RegistrationForm()
    if form.validate_on_submit():
        hashed_pw = bcrypt.generate_password_hash(form.password.data).decode('utf-8')
        user = User(username=form.username.data, email=form.email.data, password=hashed_pw)
        db.session.add(user)
        db.session.commit()
        flash('Your account has been created! You can now log in.', 'success')
        return redirect(url_for('login'))
    return render_template('register.html', form=form)

# ✅ User Login Route
@app.route("/login", methods=['GET', 'POST'])
def login():
    if current_user.is_authenticated:
        return redirect(url_for('dashboard'))
    from forms import LoginForm  # ✅ Import inside function to avoid circular import
    form = LoginForm()
    if form.validate_on_submit():
        user = User.query.filter_by(email=form.email.data).first()
        if user and bcrypt.check_password_hash(user.password, form.password.data):
            login_user(user, remember=form.remember.data)
            return redirect(url_for('dashboard'))
        flash('Login unsuccessful. Please check email and password', 'danger')
    return render_template('login.html', form=form)

# ✅ User Logout Route
@app.route("/logout")
@login_required
def logout():
    logout_user()
    return redirect(url_for('home'))

# ✅ User Dashboard
@app.route("/dashboard")
@login_required
def dashboard():
    enrollments = Enrollment.query.filter_by(user_id=current_user.id).all()
    profile_pic = url_for('static', filename='profile_pics/default_profile.jpg')
    return render_template('dashboard.html', user=current_user, profile_pic=profile_pic, enrollments=enrollments)

from upload_profile import upload_profile_bp
app.register_blueprint(upload_profile_bp)


@app.route("/explore", methods=['GET', 'POST'])
def explore():
    categories = ["All", "Data Science", "AI & ML", "Web Development", "Mobile Apps"]
    selected_category = request.form.get("category", "All")

    # Dummy course data (Replace with actual API/database query)
    all_courses = [
        {"title": "Python for Beginners", "description": "Learn Python from scratch", "image_url": "/static/python.jpg", "url": "#", "category": "Data Science"},
        {"title": "Machine Learning Basics", "description": "Introduction to ML concepts", "image_url": "/static/ml.jpg", "url": "#", "category": "AI & ML"},
        {"title": "Flutter App Development", "description": "Build mobile apps with Flutter", "image_url": "/static/flutter.jpg", "url": "#", "category": "Mobile Apps"},
    ]

    # Filter courses based on selected category
    if selected_category != "All":
        courses = [course for course in all_courses if course["category"] == selected_category]
    else:
        courses = all_courses

    return render_template("explore.html", courses=courses, categories=categories, selected_category=selected_category)


# ✅ Admin Routes for Course Management
@app.route("/admin/courses")
@login_required
@admin_required
def admin_courses():
    courses = Course.query.all()
    return render_template('courses.html', courses=courses)

@app.route("/admin/course/new", methods=['GET', 'POST'])
@login_required
@admin_required
def new_course():
    form = CourseForm()
    if form.validate_on_submit():
        course = Course(title=form.title.data, description=form.description.data)
        db.session.add(course)
        db.session.commit()
        flash("Course created!", "success")
        return redirect(url_for('admin_courses'))
    return render_template('course_form.html', form=form, legend="New Course")

@app.route("/admin/course/<int:course_id>/edit", methods=['GET', 'POST'])
@login_required
@admin_required
def edit_course(course_id):
    course = Course.query.get_or_404(course_id)
    form = CourseForm()
    if form.validate_on_submit():
        course.title = form.title.data
        course.description = form.description.data
        db.session.commit()
        flash("Course updated!", "success")
        return redirect(url_for('admin_courses'))
    elif request.method == 'GET':
        form.title.data = course.title
        form.description.data = course.description
    return render_template('course_form.html', form=form, legend="Edit Course")

@app.route("/admin/course/<int:course_id>/delete", methods=['POST'])
@login_required
@admin_required
def delete_course(course_id):
    course = Course.query.get_or_404(course_id)
    db.session.delete(course)
    db.session.commit()
    flash("Course deleted!", "success")
    return redirect(url_for('admin_courses'))
# ✅ AI Concept Extraction with Llama 3 (via Ollama)
# ✅ AI Concept Extraction with Llama 3 (via Ollama)


from flask import Flask, request, jsonify
import fitz  # PyMuPDF for PDFs
from pptx import Presentation  # PowerPoint files
from docx import Document  # Word files
import os
import json
import logging
import ollama  # AI Model for concept extraction



UPLOAD_FOLDER = "uploads"
ALLOWED_EXTENSIONS = {'pdf', 'pptx', 'docx'}
app.config['UPLOAD_FOLDER'] = UPLOAD_FOLDER

if not os.path.exists(UPLOAD_FOLDER):
    os.makedirs(UPLOAD_FOLDER)

logging.basicConfig(level=logging.INFO)

@app.route("/upload_file", methods=["POST"])
def upload_file():
    if "file" not in request.files:
        return jsonify({"error": "❌ No file uploaded"}), 400

    file = request.files["file"]
    if file.filename == "":
        return jsonify({"error": "❌ No selected file"}), 400

    file_extension = file.filename.rsplit(".", 1)[-1].lower()
    if file_extension not in ALLOWED_EXTENSIONS:
        return jsonify({"error": "❌ Unsupported file format"}), 400

    # ✅ Print file details for debugging
    print(f"📂 File Received: {file.filename}")
    print(f"📂 Saving to: {UPLOAD_FOLDER}/{file.filename}")

    filepath = os.path.join(app.config['UPLOAD_FOLDER'], file.filename)
    
    try:
        file.save(filepath)  # ✅ Save the file
        print(f"✅ File Saved: {filepath}")  # Debugging log
    except Exception as e:
        print(f"❌ Error Saving File: {e}")  # Log errors
        return jsonify({"error": f"❌ Failed to save file: {str(e)}"}), 500

    extracted_text = extract_text_from_file(filepath, file_extension)
    os.remove(filepath)  # ✅ Delete file after processing

    return jsonify({"extracted_text": extracted_text})

# 🔹 Function to Extract Text from Files
def extract_text_from_file(filepath, file_extension):
    extracted_text = ""

    try:
        if file_extension == "pdf":
            doc = fitz.open(filepath)
            for page in doc:
                extracted_text += page.get_text("text") + "\n"

        elif file_extension == "pptx":
            prs = Presentation(filepath)
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        extracted_text += shape.text + "\n"

        elif file_extension == "docx":
            doc = Document(filepath)
            for para in doc.paragraphs:
                extracted_text += para.text + "\n"

        return extracted_text.strip()

    except Exception as e:
        return str(e)


# 🔹 Function to Extract Concepts using AI
def extract_concepts_ai(text):
    if not text.strip():
        return {"error": "❌ Input text cannot be empty"}

    prompt = f"""
    Extract key concepts from the given text and categorize them as:
    - Main Concept
    - Sub-Concept
    - Related Concept

    Respond in JSON format:
    {{
      "concepts": [
        {{"id": "c1", "name": "Machine Learning", "type": "Main Concept"}},
        {{"id": "c2", "name": "Neural Networks", "type": "Sub-Concept"}},
        {{"id": "c3", "name": "Deep Learning", "type": "Related Concept"}}
      ]
    }}

    Text: {text}
    """

    try:
        response = ollama.chat(model="llama3", messages=[{"role": "user", "content": prompt}])

        # 🔹 Log Raw Response
        logging.info(f"📝 Ollama Raw Response: {response}")

        # 🔹 Validate response structure
        if not response or "message" not in response or "content" not in response["message"]:
            return {"error": "❌ No valid response from Ollama"}

        raw_response = response["message"]["content"].strip()

        # 🔹 Extract JSON part from response (if Ollama returns extra text)
        start_index = raw_response.find("{")
        end_index = raw_response.rfind("}")
        if start_index != -1 and end_index != -1:
            raw_response = raw_response[start_index:end_index+1]

        # 🔹 Parse JSON
        parsed_response = json.loads(raw_response)
        logging.info(f"✅ Parsed Response: {parsed_response}")

        return parsed_response if "concepts" in parsed_response else {"error": "❌ No concepts found. Try different text."}

    except json.JSONDecodeError:
        logging.error(f"❌ Invalid JSON from Ollama: {raw_response}")
        return {"error": "❌ AI response is not valid JSON"}

    except Exception as e:
        logging.error(f"❌ Ollama API Error: {str(e)}")
        return {"error": f"Ollama API Error: {str(e)}"}


# 🔹 Unified API for Extracting Concepts (From File OR Text)
@app.route("/extract_concepts", methods=["POST"])
def extract_concept_page():
    try:
        text = None

        # 🔹 Check if a file is uploaded
        if "file" in request.files:
            file = request.files["file"]
            if file.filename == "":
                return jsonify({"error": "❌ No selected file"}), 400

            file_extension = file.filename.rsplit(".", 1)[1].lower()
            if file_extension not in ALLOWED_EXTENSIONS:
                return jsonify({"error": "❌ Unsupported file format"}), 400

            filepath = os.path.join(app.config['UPLOAD_FOLDER'], file.filename)
            file.save(filepath)

            text = extract_text_from_file(filepath, file_extension)
            os.remove(filepath)  # Delete file after processing

        # 🔹 Check if text input is provided
        elif request.is_json:
            data = request.get_json()
            if "text" in data:
                text = data["text"].strip()

        # 🔹 Validate text extraction
        if not text or text.strip() == "":
            return jsonify({"error": "❌ No valid text found for concept extraction."}), 400

        # 🔹 Extract Concepts using AI
        concepts = extract_concepts_ai(text)
        return jsonify(concepts)

    except Exception as e:
        return jsonify({"error": f"❌ Server Error: {str(e)}"}), 500





from flask import Flask, request, jsonify
import logging
import re
import json
import logging
import random
import re
import ollama

def extract_json(text):
    """Extract a valid JSON object from AI response."""
    try:
        json_start = text.index("{")
        json_end = text.rindex("}") + 1
        json_str = text[json_start:json_end]
        return json.loads(json_str)  # ✅ Directly parse JSON safely
    except (ValueError, json.JSONDecodeError):
        logging.error(f"❌ Failed to extract valid JSON from: {text}")
        return None

def generate_quiz(concepts, difficulty="medium"):
    if not concepts:
        logging.error("❌ No concepts provided for quiz generation")
        return {"error": "❌ No concepts provided"}

    quiz = []
    question_types = ["MCQ", "True/False", "Fill in the Blank"]

    for concept in concepts:
        if isinstance(concept, str):
            concept_name = concept  
            concept_type = "General"
        elif isinstance(concept, dict):
            concept_name = concept.get("name", "Unknown Concept")
            concept_type = concept.get("type", "General")
        else:
            logging.error(f"❌ Invalid concept format: {concept}")
            continue

        question_type = random.choice(question_types)

        prompt = f"""
        You are a quiz generator. Create a **{question_type}** question for the concept "{concept_name}" ({concept_type}) with {difficulty} difficulty.
        - If MCQ, provide exactly **four** answer choices.
        
        - Always return a response in strict JSON format.

        Response Format:
        {{
          "question": "Write a question about {concept_name}.",
          "type": "{question_type}",
          "choices": ["Option A", "Option B", "Option C", "Option D"] if "{question_type}" == "MCQ" else [],
          "correct_answer": "Option B" if "{question_type}" == "MCQ" else "Correct answer text",
          "explanation": "A short explanation about why the answer is correct."
        }}
        """

        try:
            response = ollama.chat(model="llama3", messages=[{"role": "user", "content": prompt}])
            raw_response = response["message"]["content"].strip()

            # ✅ Extract JSON properly
            question_data = extract_json(raw_response)

            if not question_data:
                logging.error(f"❌ AI generated an invalid JSON response: {raw_response}")
                quiz.append({"error": "❌ AI response contained invalid JSON", "concept": concept_name})
                continue

            # ✅ Ensure MCQs have exactly 4 choices
            if question_type == "MCQ":
                if "choices" not in question_data or not isinstance(question_data["choices"], list) or len(question_data["choices"]) != 4:
                    logging.error(f"❌ MCQ question missing proper options: {question_data}")
                    quiz.append({"error": "❌ MCQ options were not properly generated", "concept": concept_name})
                    continue

            # ✅ Ensure `choices` exists in True/False & Fill-in-the-Blank
            if question_type in ["True/False", "Fill in the Blank"] and "choices" not in question_data:
                question_data["choices"] = []

            question_data["difficulty"] = difficulty
            quiz.append(question_data)

        except Exception as e:
            logging.error(f"❌ ERROR in generate_quiz: {e}")
            quiz.append({"error": "❌ Failed to generate question", "concept": concept_name})

    return quiz

import traceback

@app.route("/generate_quiz", methods=["POST"])
def generate_quiz_api():
    try:
        data = request.json
        logging.info(f"📩 Received data: {data}")

        if not data or "concepts" not in data:
            logging.error("❌ Invalid request: 'concepts' field is required")
            return jsonify({"error": "❌ 'concepts' field is required", "quiz": []}), 400

        concepts = data["concepts"]

        # ✅ Allow both strings and dictionaries as concepts
        if not isinstance(concepts, list) or not all(isinstance(c, (str, dict)) for c in concepts):
            logging.error("❌ 'concepts' should be a list of strings or dictionaries")
            return jsonify({"error": "❌ 'concepts' should be a list of strings or dictionaries", "quiz": []}), 400

        difficulty = data.get("difficulty", "medium")
        logging.info(f"📚 Concepts: {concepts}, Difficulty: {difficulty}")

        quiz = generate_quiz(concepts, difficulty)
        logging.info(f"✅ Quiz generated successfully with {len(quiz)} questions")

        return jsonify({"quiz": quiz})  # ✅ Send quiz data as JSON

    except Exception as e:
        logging.error(f"❌ Error in /generate_quiz: {str(e)}")
        logging.error(traceback.format_exc())  # ✅ Print full error traceback
        return jsonify({"error": f"❌ Server Error: {str(e)}", "quiz": []}), 500




# ✅ Function: Evaluate Answers
def evaluate_answers(quiz, user_answers):
    if not quiz or not user_answers:
        return {"error": "❌ Quiz or answers missing"}

    evaluation = []
    for i, question in enumerate(quiz):
        correct_answer = str(question.get("correct_answer", "")).strip()
        user_response = str(user_answers[i]).strip()

        prompt = f"""
        Question: "{question['question']}"
        Correct Answer: "{correct_answer}"
        User Answer: "{user_response}"

        Evaluate and provide:
        - Score (0-10)
        - Explanation
        - Improvement tips

        Respond in JSON format:
        {{
          "score": 8,
          "feedback": "Your explanation",
          "suggestions": "How to improve?"
        }}
        """

        try:
            response = ollama.chat(model="llama3", messages=[{"role": "user", "content": prompt}])
            raw_response = response["message"]["content"].strip()
            evaluation.append(json.loads(raw_response))

        except Exception as e:
            logging.error(f"❌ ERROR in evaluate_answers: {e}")
            evaluation.append({"error": "Failed to evaluate answer"})

    return evaluation


# ✅ API Endpoint: Evaluate Answers
@app.route('/evaluate_answers', methods=['POST'])
def evaluate_answers_api():
    try:
        data = request.get_json()
        logging.info(f"📩 Received Evaluation Request: {data}")

        quiz = data.get("quiz", [])
        user_answers = data.get("answers", [])

        if not quiz or not user_answers:
            return jsonify({"error": "❌ Quiz or answers missing"}), 400

        evaluation = evaluate_answers(quiz, user_answers)
        return jsonify({"evaluation": evaluation})

    except Exception as e:
        logging.error(f"❌ Error in /evaluate_answers: {str(e)}")
        return jsonify({"error": f"❌ Server Error: {str(e)}"}), 500

# ✅ Run Flask App
if __name__ == '__main__':
    app.run(debug=True)